"""Generic document processing service supporting multiple file formats."""
import logging
from typing import List, Dict, Any
from pathlib import Path
from dataclasses import dataclass

logger = logging.getLogger(__name__)


@dataclass
class TextSegment:
    """Text segment with position information."""
    text: str
    page_number: int
    char_start: int  # Character position within the document
    char_end: int
    bbox: tuple  # Bounding box (x0, y0, x1, y1) for potential future use


class DocumentProcessor:
    """Process various document formats and extract text with position tracking."""

    SUPPORTED_EXTENSIONS = {
        '.pdf': 'PDF Document',
        '.docx': 'Word Document',
        '.doc': 'Word Document (Legacy)',
        '.txt': 'Text File',
        '.md': 'Markdown File',
        '.pptx': 'PowerPoint Presentation',
        '.xlsx': 'Excel Spreadsheet',
        '.csv': 'CSV File',
        '.html': 'HTML Document',
        '.htm': 'HTML Document',
        '.rtf': 'Rich Text Format',
        '.odt': 'OpenDocument Text',
    }

    def __init__(self):
        """Initialize document processor."""
        pass

    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions."""
        return list(self.SUPPORTED_EXTENSIONS.keys())

    def is_supported(self, file_path: str) -> bool:
        """Check if file format is supported."""
        ext = Path(file_path).suffix.lower()
        return ext in self.SUPPORTED_EXTENSIONS

    def extract_text_with_positions(self, file_path: str) -> tuple[List[TextSegment], Dict[str, Any]]:
        """
        Extract text from any supported document format with position tracking.

        Args:
            file_path: Path to the document file

        Returns:
            Tuple of (list of TextSegments, document metadata)
        """
        ext = Path(file_path).suffix.lower()

        if ext == '.pdf':
            return self._extract_from_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return self._extract_from_docx(file_path)
        elif ext == '.txt':
            return self._extract_from_txt(file_path)
        elif ext == '.md':
            return self._extract_from_markdown(file_path)
        elif ext == '.pptx':
            return self._extract_from_pptx(file_path)
        elif ext in ['.xlsx', '.csv']:
            return self._extract_from_excel(file_path)
        elif ext in ['.html', '.htm']:
            return self._extract_from_html(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def _extract_from_pdf(self, pdf_path: str) -> tuple[List[TextSegment], Dict[str, Any]]:
        """Extract text from PDF with position tracking."""
        import pdfplumber

        segments = []
        metadata = {
            "page_count": 0,
            "total_chars": 0,
            "file_size": Path(pdf_path).stat().st_size,
            "file_type": "pdf"
        }

        try:
            with pdfplumber.open(pdf_path) as pdf:
                metadata["page_count"] = len(pdf.pages)
                document_char_offset = 0

                for page_num, page in enumerate(pdf.pages, start=1):
                    page_text = page.extract_text()

                    if not page_text:
                        logger.warning(f"No text found on page {page_num}")
                        continue

                    # Create segments based on paragraphs
                    paragraphs = page_text.split('\n\n')
                    page_char_offset = 0

                    for para in paragraphs:
                        if para.strip():
                            segment = TextSegment(
                                text=para.strip(),
                                page_number=page_num,
                                char_start=document_char_offset + page_char_offset,
                                char_end=document_char_offset + page_char_offset + len(para),
                                bbox=(0, 0, page.width, page.height)
                            )
                            segments.append(segment)

                        page_char_offset += len(para) + 2  # +2 for '\n\n'

                    metadata["total_chars"] += len(page_text)
                    document_char_offset += len(page_text)

                logger.info(f"Extracted {len(segments)} segments from {metadata['page_count']} pages (PDF)")

        except Exception as e:
            logger.error(f"Error processing PDF {pdf_path}: {e}")
            raise

        return segments, metadata

    def _extract_from_docx(self, docx_path: str) -> tuple[List[TextSegment], Dict[str, Any]]:
        """Extract text from DOCX file."""
        from docx import Document

        segments = []
        metadata = {
            "page_count": 1,  # DOCX doesn't have explicit page concept
            "total_chars": 0,
            "file_size": Path(docx_path).stat().st_size,
            "file_type": "docx"
        }

        try:
            doc = Document(docx_path)
            document_char_offset = 0

            for para_idx, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if text:
                    segment = TextSegment(
                        text=text,
                        page_number=1,  # All on "page 1" for DOCX
                        char_start=document_char_offset,
                        char_end=document_char_offset + len(text),
                        bbox=(0, 0, 0, 0)
                    )
                    segments.append(segment)
                    document_char_offset += len(text) + 2

            metadata["total_chars"] = document_char_offset
            logger.info(f"Extracted {len(segments)} paragraphs from DOCX")

        except Exception as e:
            logger.error(f"Error processing DOCX {docx_path}: {e}")
            raise

        return segments, metadata

    def _extract_from_txt(self, txt_path: str) -> tuple[List[TextSegment], Dict[str, Any]]:
        """Extract text from plain text file."""
        segments = []
        metadata = {
            "page_count": 1,
            "total_chars": 0,
            "file_size": Path(txt_path).stat().st_size,
            "file_type": "txt"
        }

        try:
            with open(txt_path, 'r', encoding='utf-8') as f:
                content = f.read()

            # Split by double newlines (paragraphs)
            paragraphs = content.split('\n\n')
            document_char_offset = 0

            for para in paragraphs:
                text = para.strip()
                if text:
                    segment = TextSegment(
                        text=text,
                        page_number=1,
                        char_start=document_char_offset,
                        char_end=document_char_offset + len(text),
                        bbox=(0, 0, 0, 0)
                    )
                    segments.append(segment)
                    document_char_offset += len(para) + 2

            metadata["total_chars"] = len(content)
            logger.info(f"Extracted {len(segments)} paragraphs from TXT")

        except Exception as e:
            logger.error(f"Error processing TXT {txt_path}: {e}")
            raise

        return segments, metadata

    def _extract_from_markdown(self, md_path: str) -> tuple[List[TextSegment], Dict[str, Any]]:
        """Extract text from Markdown file."""
        import markdown
        from html.parser import HTMLParser

        class HTMLTextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text_parts = []

            def handle_data(self, data):
                self.text_parts.append(data)

        segments = []
        metadata = {
            "page_count": 1,
            "total_chars": 0,
            "file_size": Path(md_path).stat().st_size,
            "file_type": "markdown"
        }

        try:
            with open(md_path, 'r', encoding='utf-8') as f:
                md_content = f.read()

            # Convert markdown to HTML, then extract plain text
            html = markdown.markdown(md_content)
            extractor = HTMLTextExtractor()
            extractor.feed(html)
            text = ' '.join(extractor.text_parts)

            # Also use raw markdown split by sections
            paragraphs = md_content.split('\n\n')
            document_char_offset = 0

            for para in paragraphs:
                text = para.strip()
                if text:
                    segment = TextSegment(
                        text=text,
                        page_number=1,
                        char_start=document_char_offset,
                        char_end=document_char_offset + len(text),
                        bbox=(0, 0, 0, 0)
                    )
                    segments.append(segment)
                    document_char_offset += len(para) + 2

            metadata["total_chars"] = len(md_content)
            logger.info(f"Extracted {len(segments)} sections from Markdown")

        except Exception as e:
            logger.error(f"Error processing Markdown {md_path}: {e}")
            raise

        return segments, metadata

    def _extract_from_pptx(self, pptx_path: str) -> tuple[List[TextSegment], Dict[str, Any]]:
        """Extract text from PowerPoint file."""
        from pptx import Presentation

        segments = []
        metadata = {
            "page_count": 0,
            "total_chars": 0,
            "file_size": Path(pptx_path).stat().st_size,
            "file_type": "pptx"
        }

        try:
            prs = Presentation(pptx_path)
            metadata["page_count"] = len(prs.slides)
            document_char_offset = 0

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []

                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)

                # Combine all text from the slide
                combined_text = '\n'.join(slide_text).strip()

                if combined_text:
                    segment = TextSegment(
                        text=combined_text,
                        page_number=slide_num,
                        char_start=document_char_offset,
                        char_end=document_char_offset + len(combined_text),
                        bbox=(0, 0, 0, 0)
                    )
                    segments.append(segment)
                    document_char_offset += len(combined_text) + 2

            metadata["total_chars"] = document_char_offset
            logger.info(f"Extracted {len(segments)} slides from PPTX")

        except Exception as e:
            logger.error(f"Error processing PPTX {pptx_path}: {e}")
            raise

        return segments, metadata

    def _extract_from_excel(self, excel_path: str) -> tuple[List[TextSegment], Dict[str, Any]]:
        """Extract text from Excel file."""
        import openpyxl
        import csv

        segments = []
        ext = Path(excel_path).suffix.lower()

        metadata = {
            "page_count": 0,
            "total_chars": 0,
            "file_size": Path(excel_path).stat().st_size,
            "file_type": ext[1:]  # Remove the dot
        }

        try:
            document_char_offset = 0

            if ext == '.xlsx':
                wb = openpyxl.load_workbook(excel_path, read_only=True)
                metadata["page_count"] = len(wb.sheetnames)

                for sheet_num, sheet_name in enumerate(wb.sheetnames, start=1):
                    sheet = wb[sheet_name]
                    rows = []

                    for row in sheet.iter_rows(values_only=True):
                        row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                        if row_text.strip():
                            rows.append(row_text)

                    if rows:
                        combined_text = '\n'.join(rows)
                        segment = TextSegment(
                            text=f"Sheet: {sheet_name}\n{combined_text}",
                            page_number=sheet_num,
                            char_start=document_char_offset,
                            char_end=document_char_offset + len(combined_text),
                            bbox=(0, 0, 0, 0)
                        )
                        segments.append(segment)
                        document_char_offset += len(combined_text) + 2

            elif ext == '.csv':
                metadata["page_count"] = 1
                with open(excel_path, 'r', encoding='utf-8') as f:
                    reader = csv.reader(f)
                    rows = []
                    for row in reader:
                        row_text = ' | '.join(row)
                        if row_text.strip():
                            rows.append(row_text)

                    if rows:
                        combined_text = '\n'.join(rows)
                        segment = TextSegment(
                            text=combined_text,
                            page_number=1,
                            char_start=0,
                            char_end=len(combined_text),
                            bbox=(0, 0, 0, 0)
                        )
                        segments.append(segment)

            metadata["total_chars"] = document_char_offset
            logger.info(f"Extracted {len(segments)} sheets/rows from {ext.upper()}")

        except Exception as e:
            logger.error(f"Error processing {ext.upper()} {excel_path}: {e}")
            raise

        return segments, metadata

    def _extract_from_html(self, html_path: str) -> tuple[List[TextSegment], Dict[str, Any]]:
        """Extract text from HTML file."""
        from html.parser import HTMLParser

        class HTMLTextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text_parts = []
                self.in_script_or_style = False

            def handle_starttag(self, tag, attrs):
                if tag in ['script', 'style']:
                    self.in_script_or_style = True

            def handle_endtag(self, tag):
                if tag in ['script', 'style']:
                    self.in_script_or_style = False

            def handle_data(self, data):
                if not self.in_script_or_style:
                    text = data.strip()
                    if text:
                        self.text_parts.append(text)

        segments = []
        metadata = {
            "page_count": 1,
            "total_chars": 0,
            "file_size": Path(html_path).stat().st_size,
            "file_type": "html"
        }

        try:
            with open(html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            extractor = HTMLTextExtractor()
            extractor.feed(html_content)

            document_char_offset = 0
            for text_part in extractor.text_parts:
                if text_part:
                    segment = TextSegment(
                        text=text_part,
                        page_number=1,
                        char_start=document_char_offset,
                        char_end=document_char_offset + len(text_part),
                        bbox=(0, 0, 0, 0)
                    )
                    segments.append(segment)
                    document_char_offset += len(text_part) + 2

            metadata["total_chars"] = document_char_offset
            logger.info(f"Extracted {len(segments)} text blocks from HTML")

        except Exception as e:
            logger.error(f"Error processing HTML {html_path}: {e}")
            raise

        return segments, metadata

    def validate_document(self, file_path: str) -> bool:
        """
        Validate that the file is a valid document.

        Args:
            file_path: Path to the document file

        Returns:
            True if valid, False otherwise
        """
        try:
            if not self.is_supported(file_path):
                return False

            # Try to extract text to verify it's valid
            segments, _ = self.extract_text_with_positions(file_path)
            return len(segments) > 0

        except Exception as e:
            logger.error(f"Document validation failed for {file_path}: {e}")
            return False
